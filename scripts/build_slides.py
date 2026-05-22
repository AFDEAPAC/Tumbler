#!/usr/bin/env python3
"""
Build tumbler-slides-gtac2026.pptx from the HIPER template by
python-pptx.  Strategy: open template, strip its 21 sample slides while
keeping the slide masters + layouts + theme intact, then add 13 new
slides using the same layouts.
"""
import os
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from pathlib import Path as _Path
REPO = _Path(__file__).resolve().parents[1]



A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def _a(tag):
    return f'{{{A_NS}}}{tag}'


def _force_cell_borders(cell, *, width_emu=12700, hex_rgb='000000'):
    """Apply solid borders on all 4 edges of a pptx table cell.
    width_emu defaults to 1pt (12700 EMU = 1 pt).  Works around the
    fact that python-pptx tables inherit master-style borders that
    LibreOffice often renders as invisible."""
    tcPr = cell._tc.get_or_add_tcPr()
    for edge in ('lnL', 'lnR', 'lnT', 'lnB'):
        # remove any pre-existing edge so style inheritance is overridden
        for old in tcPr.findall(_a(edge)):
            tcPr.remove(old)
        ln = etree.SubElement(tcPr, _a(edge))
        ln.set('w', str(width_emu))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        solidFill = etree.SubElement(ln, _a('solidFill'))
        srgb = etree.SubElement(solidFill, _a('srgbClr'))
        srgb.set('val', hex_rgb)
        prstDash = etree.SubElement(ln, _a('prstDash'))
        prstDash.set('val', 'solid')

TEMPLATE = str(REPO / 'docs/templates/GTAC2025-HIPER.pptx')
OUT      = str(REPO / 'docs/slides/tumbler-slides-gtac2026.pptx')
FIG_DIR  = REPO / 'docs/paper/figures'

prs = Presentation(TEMPLATE)


# --- Patch hardcoded 'GTAC 2025' year strings in slide master + layouts ---
# The HIPER template ships with the previous year baked into:
#   * slideMaster1.xml:   '|    AMD GTAC 2025' footer (shown on every slide)
#   * slideLayout1.xml:   'GTAC 2025'           subtitle on the Title Slide
# Walk every <a:t> text run in masters and layouts and rewrite 2025 -> 2026
# so the conference year matches the current submission.
def _patch_year_in_part(part, old='2025', new='2026'):
    if part is None or part.element is None:
        return 0
    n = 0
    for t in part.element.iter(qn('a:t')):
        if t.text and old in t.text:
            t.text = t.text.replace(old, new)
            n += 1
    return n

for master in prs.slide_masters:
    _patch_year_in_part(master)
    for layout in master.slide_layouts:
        _patch_year_in_part(layout)


# --- strip all existing slides (keep masters + layouts) ---
sld_id_lst = prs.slides._sldIdLst
for sld in list(sld_id_lst):
    rId = sld.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass
    sld_id_lst.remove(sld)


# ----- layout lookup -----
LAY = {l.name: l for l in prs.slide_masters[0].slide_layouts}
TITLE_FULL  = LAY['Title Slide - Full Width']
TITLE_BODY  = LAY['Title and Content']
TITLE_ONLY  = LAY['Title Only']
TITLE_SIDE  = LAY['Title and text side by side 2']
BLANK       = LAY['Blank']
CLOSE_AMD   = LAY['2_AMD Logo']

SLIDE_W = prs.slide_width   # 12192000 = 13.33"
SLIDE_H = prs.slide_height  # 6858000  = 7.50"


def set_text(tf, text, *, font_size=None, bold=None, align=None):
    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    if font_size is not None:
        r.font.size = Pt(font_size)
    if bold is not None:
        r.font.bold = bold


def add_bullets(tf, items, *, font_size=18):
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = item
        r.font.size = Pt(font_size)


def find_placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def find_title(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.placeholder_format is not None:
            if sh.placeholder_format.idx == 0:
                return sh
    return None


def add_centered_image(slide, image_path, *, top=Inches(1.4),
                       max_h=Inches(5.6), max_w=Inches(11.5)):
    """Insert an image centred horizontally, scaled to fit within max_w/max_h."""
    from PIL import Image as PILImage
    with PILImage.open(image_path) as im:
        w_px, h_px = im.size
    aspect = w_px / h_px
    cur_w = max_w
    cur_h = int(cur_w / aspect)
    if cur_h > max_h:
        cur_h = max_h
        cur_w = int(cur_h * aspect)
    left = int((SLIDE_W - cur_w) // 2)
    slide.shapes.add_picture(str(image_path), left, top,
                             width=cur_w, height=cur_h)


def add_full_bleed_image(slide, image_path):
    slide.shapes.add_picture(str(image_path), 0, 0,
                             width=SLIDE_W, height=SLIDE_H)


def add_text_box(slide, left, top, width, height, text, *,
                 font_size=18, bold=False, align=PP_ALIGN.LEFT,
                 colour=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    if colour is not None:
        r.font.color.rgb = colour
    return tb


def add_table(slide, headers, rows, left, top, width, height,
              header_fs=14, body_fs=12):
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                         width, height)
    tbl = table_shape.table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = ''
        tf = cell.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.size = Pt(header_fs)
        r.font.bold = True
        _force_cell_borders(cell)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(1 + i, j)
            cell.text = ''
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = val
            r.font.size = Pt(body_fs)
            _force_cell_borders(cell)
    return tbl


# =====================================================================
#                              SLIDES
# =====================================================================

# ---- Slide 1: Title ----
s = prs.slides.add_slide(TITLE_FULL)
title_sh = find_title(s)
if title_sh is not None:
    # Short slide title — the long paper title overlaps the layout's
    # decorative AMD background text.  The paper docx keeps the full
    # title; the deck uses a shorter version.
    set_text(title_sh.text_frame,
             'Tumbler — Surviving GPU Faults in Production AI',
             font_size=40, bold=True)
sub = find_placeholder(s, 12)
if sub is not None:
    tf = sub.text_frame
    tf.clear()
    # 'GTAC 2026' is provided by the slide layout's orange subtitle text,
    # so we do not repeat it here.
    lines = [
        'Three-tier  ·  Unified  ·  Multi-tenant  ·  Bounded-wait  ·  '
        'Layered  ·  Escape-hatch  ·  Runtime',
        '',
        'Chun-Hung Wang*  /  Clement Lin  /  Jeremy Liao',
        'AMD',
        '',
        '* Chun-Hung.Wang@amd.com (corresponding)',
    ]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(18) if i == 0 else Pt(14)
        if i == 0:
            r.font.bold = True


# ---- Slide 2: Hook — cockpit frozen ----
s = prs.slides.add_slide(BLANK)
add_full_bleed_image(s, FIG_DIR / 'fig-02-car-cockpit-frozen.png')
add_text_box(s, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.6),
             "Picture this — but on production AI",
             font_size=28, bold=True, colour=RGBColor(0xFF, 0xFF, 0xFF))
add_text_box(s, Inches(0.6), Inches(6.4), Inches(12.0), Inches(0.7),
             "If the touchscreen GPU faults on the highway, you want one "
             "frozen frame — not the whole infotainment process dead.",
             font_size=18, colour=RGBColor(0xFF, 0xFF, 0xFF))


# ---- Slide 3: Why AI workloads are different ----
s = prs.slides.add_slide(TITLE_BODY)
set_text(find_title(s).text_frame,
         "Why AI workloads are different from HPC",
         font_size=32, bold=True)
body = find_placeholder(s, 10)
add_bullets(body.text_frame, [
    "HPC: 1 host process owns 1 GPU. A VM fault means data corruption — "
    "the only safe action is to abort.",
    "AI serving: 1 host process owns 100s of concurrent tenant requests "
    "sharing the same GPUs, SDMA engines, signals, and KFD pinned memory.",
    "Same fault, very different cost: one bad request vs one dead serving "
    "node taking every co-tenant down with it.",
    "ROCm was built for the first case. Production AI inverts the "
    "assumption — and needs an opt-in escape hatch.",
], font_size=20)


# ---- Slide 4: Today — abort-fatal default ----
s = prs.slides.add_slide(TITLE_ONLY)
set_text(find_title(s).text_frame,
         "Today — one fault aborts the whole serving process",
         font_size=30, bold=True)
add_centered_image(s, FIG_DIR / 'fig-03-ai-workload-fault-propagation.png',
                   top=Inches(1.3), max_h=Inches(5.6), max_w=Inches(10.5))


# ---- Slide 5: Why ROCm has no firewall ----
s = prs.slides.add_slide(TITLE_BODY)
set_text(find_title(s).text_frame,
         "Why ROCm has no built-in survival firewall",
         font_size=30, bold=True)
body = find_placeholder(s, 10)
add_bullets(body.text_frame, [
    "ROCr inherits the HSA spec: one logical process owns the GPU. "
    "Three abort sites in runtime.cpp (memory error / HW exception / VM "
    "fault) all call std::abort() unconditionally.",
    "CLR WaitForSignal() has a 4-sec inner poll, but the outer while loop "
    "is unbounded — a wedged signal parks the caller thread forever.",
    "amdgpu waits with MAX_SCHEDULE_TIMEOUT in kcl_drm_suballoc_new() and "
    "on a global mutex in the GTT-window path.",
    "HSA_SIGNAL_WAIT_ABORT_TIMEOUT exists — but it also aborts the "
    "process at the deadline. Right for batch HPC, wrong for AI serving.",
], font_size=18)


# ---- Slide 6: Tumbler concept (with acronym reveal) ----
s = prs.slides.add_slide(TITLE_ONLY)
set_text(find_title(s).text_frame,
         "TUMBLER — a bounded-wait firewall that stops the blast radius",
         font_size=26, bold=True)
# Subtitle line with the acronym expansion, styled smaller + lighter.
add_text_box(s, Inches(0.6), Inches(0.9), Inches(12.0), Inches(0.5),
             "Three-tier  ·  Unified  ·  Multi-tenant  ·  Bounded-wait  ·  "
             "Layered  ·  Escape-hatch  ·  Runtime",
             font_size=14, colour=RGBColor(0xAA, 0xC8, 0xEE))
add_centered_image(s, FIG_DIR / 'fig-04-tumbler-firewall-layered.png',
                   top=Inches(1.6), max_h=Inches(5.3), max_w=Inches(10.5))


# ---- Slide 7: Architecture ----
s = prs.slides.add_slide(TITLE_ONLY)
set_text(find_title(s).text_frame,
         "Architecture — knobs across the ROCm stack",
         font_size=28, bold=True)
add_centered_image(s, FIG_DIR / 'fig-07-stack-layers-knobs.png',
                   top=Inches(1.3), max_h=Inches(5.6), max_w=Inches(10.5))


# ---- Slide 8: Three layers of defense ----
s = prs.slides.add_slide(TITLE_ONLY)
set_text(find_title(s).text_frame,
         "Three layers of defense",
         font_size=32, bold=True)
add_table(
    s,
    headers=['Layer', 'Guards against', 'How it protects'],
    rows=[
        ['CLR\n(HIP host runtime)',
         'A host thread parked forever on a wedged HSA completion '
         'signal — the serving scheduler cannot reroute the affected '
         'request.',
         'Bounded outer-loop wait on WaitForSignal(). At the deadline '
         'the signal is force-completed and the caller returns. The '
         'serving framework decides whether to retry or fail the '
         'one request.'],
        ['ROCr\n(HSA runtime)',
         'std::abort() on the three runtime.cpp abort sites '
         '(memory error, non-ECC HW exception, VM fault) — one fault '
         'kills the whole process and takes every co-tenant request '
         'down with it. Plus SDMA writers spinning forever; '
         'InterruptSignal waits parked.',
         'Opt-in non-abort branches at every abort site (ECC stays '
         'abort-fatal for data integrity). Bounded caps on the '
         'signal-wait inner loop and on the SDMA yield loop. '
         'Failure surfaces as an HSA_STATUS error, never as a '
         'process abort.'],
        ['amdgpu\n(kernel driver)',
         'A wedged SDMA ring holding the global GTT-window mutex '
         'and starving every VRAM-touching ioctl on the device. '
         'A pinned BO whose owner exited without unpinning '
         'permanently holding VRAM. Unbounded dma_fence_wait inside '
         'suballoc.',
         'Bounded timeouts at every fence and lock site (suballoc, '
         'GTT window, KFD free / unpin). KFD pin-reaper subsystem '
         'harvests orphan pins on a deadline. Failure surfaces to '
         'user space as -ETIME, never as a hung ioctl.'],
    ],
    left=Inches(0.3), top=Inches(1.2),
    width=Inches(12.7), height=Inches(5.7),
    header_fs=16, body_fs=12)


# ---- Slide 9: Production result ----
s = prs.slides.add_slide(TITLE_ONLY)
set_text(find_title(s).text_frame,
         "Production result — single-digit minutes vs healthy waveform",
         font_size=26, bold=True)
add_centered_image(s, FIG_DIR / 'fig-05-latency-waveform.png',
                   top=Inches(1.3), max_h=Inches(5.6), max_w=Inches(12.5))


# ---- Slide 10: Multi-tenant survival ----
s = prs.slides.add_slide(TITLE_ONLY)
set_text(find_title(s).text_frame,
         "Multi-tenant survival — fault blast radius is one request",
         font_size=26, bold=True)
add_centered_image(s, FIG_DIR / 'fig-06-multi-tenant-survival.png',
                   top=Inches(1.3), max_h=Inches(5.6), max_w=Inches(10.5))


# ---- (Slide formerly #11 'Upstream PRs' dropped — PR details live in
#       the meta-repo / paper References; the deck does not need them.) ----


# ---- Slide 11: Future work ----
s = prs.slides.add_slide(TITLE_BODY)
set_text(find_title(s).text_frame,
         "Future work — AI-driven survival",
         font_size=30, bold=True)
body = find_placeholder(s, 10)
add_bullets(body.text_frame, [
    "Auto-tune Tumbler deadlines from observed workload signatures "
    "(serving-framework tail-latency budget, batch dwell time, in-flight "
    "request count).",
    "Reinforcement-learning-driven survival policy: learn when to fail "
    "a request gracefully vs wait out a transient hiccup.",
    "Feedback channels into the inference framework so the higher-level "
    "scheduler can adjust load on graceful-fail events.",
    "Long-tail items: KFD slow-ioctl bounded path, GTT multi-window "
    "lock sharding (deferred for future PRs).",
], font_size=20)


# ---- Slide 12: Q&A / contact ----
s = prs.slides.add_slide(TITLE_BODY)
set_text(find_title(s).text_frame,
         "Thank you — questions?",
         font_size=36, bold=True)
body = find_placeholder(s, 10)
add_bullets(body.text_frame, [
    "Chun-Hung Wang  —  Chun-Hung.Wang@amd.com  (corresponding)",
    "Clement Lin     —  Clement.Lin@amd.com",
    "Jeremy Liao     —  Jeremy.Liao@amd.com",
    "",
    "Code:   github.com/AFDEAPAC/Tumbler",
    "PRs:    ROCm/rocm-systems #6328, #6329  /  ROCm/amdgpu #214, #215, #216",
], font_size=20)


# ---- Slide 13 (optional): AMD logo closing slide ----
s = prs.slides.add_slide(CLOSE_AMD)


# --- core properties (Tumbler) -----------------------------------------
cp = prs.core_properties
cp.title = 'TUMBLER: Three-tier Bounded-Wait Survival Stack for ROCm'
cp.subject = 'AMD ROCm CLR + ROCr + amdgpu fault-tolerant runtime/driver stack'
cp.author = 'Chun-Hung Wang'
cp.last_modified_by = 'Chun-Hung Wang'
cp.keywords = 'ROCm; HIP; HSA; amdgpu; bounded-wait; survival; multi-tenant; AI serving; MI300X'
cp.comments = 'GTAC 2026 submission. Generated by scripts/build_slides.py from the AMD corporate template.'
# -----------------------------------------------------------------------
prs.save(OUT)
print(f'wrote: {OUT}')
print(f'  size: {os.path.getsize(OUT)} bytes')
print(f'  slides: {len(prs.slides)}')


# --- scrub docProps/app.xml leftover template fields (Tumbler) -------
import zipfile, shutil, re as _re, os as _os, tempfile as _tempfile
def _scrub_app_xml(path):
    tmp_fd, tmp_path = _tempfile.mkstemp(suffix='.xlsx')
    _os.close(tmp_fd)
    with zipfile.ZipFile(path, 'r') as zin, zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == 'docProps/app.xml':
                txt = data.decode('utf-8', errors='ignore')
                txt = _re.sub(r'<Manager>[^<]*</Manager>', '<Manager></Manager>', txt)
                txt = _re.sub(r'<Company>\s*</Company>', '<Company>Advanced Micro Devices</Company>', txt)
                data = txt.encode('utf-8')
            zout.writestr(item, data)
    shutil.move(tmp_path, path)
    _os.chmod(path, 0o644)

_scrub_app_xml(OUT)
print('scrubbed app.xml placeholders in', OUT)
# ---------------------------------------------------------------------
